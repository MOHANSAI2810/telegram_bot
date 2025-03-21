import os
import re
import asyncio
import requests
import google.generativeai as genai
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, CallbackContext
from dotenv import load_dotenv
import base64
import PyPDF2
from docx import Document
from pptx import Presentation
import html  # For escaping special characters

# Load environment variables
load_dotenv()

# Get API keys from environment variables
TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
GEMINI_API_KEY = os.getenv("api_key")
OPENWEATHERMAP_API_KEY = os.getenv("OPENWEATHERMAP_API_KEY")
NEWS_API_KEY = os.getenv("NEWS_API_KEY")
CLARIFAI_API_KEY = os.getenv("CLARIFAI_API_KEY")
CLARIFAI_MODEL_URL = os.getenv("CLARIFAI_MODEL_URL")

# Configure Gemini API
genai.configure(api_key=GEMINI_API_KEY)

# Define Gemini model
model = genai.GenerativeModel("gemini-1.5-flash")

# Inappropriate keywords filter
INAPPROPRIATE_KEYWORDS = ["adult", "porn", "sex", "violence", "drugs", "hate"]

# Custom responses
CUSTOM_RESPONSES = {
    "what is your name": "I am Mohan's Mini Chatbot, and I am here to help you!",
    "who are you": "I am Mohan's Mini Chatbot, and I am here to help you!",
    "who is your founder": "My boss is Mr. Mohan.",
    "who created you": "My boss is Mr. Mohan.",
    "who is your boss": "My boss is Mr. Mohan.",
    "what can you do": "I can help you with a variety of tasks, such as answering questions, analyzing files, and more!",
    "how are you": "I'm fine, thank you!",
    "what is your purpose": "My purpose is to assist you with your queries and make your life easier.",
}

def is_inappropriate(content):
    """Check if content contains inappropriate keywords."""
    content = content.lower()
    return any(keyword in content for keyword in INAPPROPRIATE_KEYWORDS)

def format_response(response_text):
    """Format the response text by making text inside ** ** bold."""
    # Escape special characters to avoid HTML parsing errors
    response_text = html.escape(response_text)
    # Replace **text** with <b>text</b>
    return re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", response_text)

def get_custom_response(user_message):
    """Check if the user message matches any custom response."""
    user_message = user_message.lower().strip()
    return CUSTOM_RESPONSES.get(user_message)

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return None

def extract_text_from_docx(file_path):
    """Extract text from a DOCX file."""
    try:
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return None

def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file."""
    try:
        ppt = Presentation(file_path)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return None

def summarize_text_with_gemini(text):
    """Summarize text using the Gemini API."""
    try:
        prompt = f"Summarize the following text and provide the key points:\n{text}"
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        print(f"Error summarizing text with Gemini: {e}")
        return None

def analyze_image_with_clarifai(image_path):
    """Send the image to Clarifai API and return the analysis results."""
    with open(image_path, "rb") as image_file:
        image_data = base64.b64encode(image_file.read()).decode("utf-8")

    headers = {
        "Authorization": f"Key {CLARIFAI_API_KEY}",
        "Content-Type": "application/json",
    }

    payload = {
        "inputs": [
            {
                "data": {
                    "image": {
                        "base64": image_data,
                    }
                }
            }
        ]
    }

    try:
        response = requests.post(CLARIFAI_MODEL_URL, headers=headers, json=payload)
        response.raise_for_status()
        results = response.json()
        concepts = results["outputs"][0]["data"]["concepts"]
        labels_with_scores = [(concept["name"], concept["value"]) for concept in concepts]
        return labels_with_scores
    except Exception as e:
        print(f"Error analyzing image with Clarifai: {e}")
        return None

def generate_paragraph_with_gemini(labels_with_scores):
    """Send the Clarifai response to Gemini API and generate a descriptive paragraph."""
    labels_text = ", ".join([f"{label} ({score * 100:.2f}%)" for label, score in labels_with_scores])
    prompt = f"Describe the following image contents in a paragraph: {labels_text}"
    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        print(f"Error generating paragraph with Gemini: {e}")
        return None

async def get_weather(city):
    """Get weather information for a city."""
    url = f"https://api.openweathermap.org/data/2.5/weather?q={city}&appid={OPENWEATHERMAP_API_KEY}&units=metric"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        temperature = data["main"]["temp"]
        weather_description = data["weather"][0]["description"]
        return f"The weather in {city} is {weather_description} with a temperature of {temperature}°C."
    else:
        return "Sorry, I couldn't fetch the weather."

async def get_news():
    """Get the latest news headlines."""
    url = f"https://newsapi.org/v2/top-headlines?country=us&apiKey={NEWS_API_KEY}"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        articles = data["articles"][:5]
        news_message = "Here are the latest headlines:\n" + "\n".join([article["title"] for article in articles])
        return news_message
    else:
        return "Sorry, I couldn't fetch the news."

async def translate_text(text, target_language):
    """Translate text to the target language."""
    url = f"https://api.mymemory.translated.net/get?q={text}&langpair=en|{target_language}"
    response = requests.get(url)
    if response.status_code == 200:
        translation_data = response.json()
        if translation_data.get("responseStatus") == 200:
            return translation_data["responseData"]["translatedText"]
    return "Sorry, I couldn't translate that."

async def send_response(update, context, text):
    """Send the response with a typing effect."""
    # Show typing indicator
    await context.bot.send_chat_action(chat_id=update.message.chat_id, action="typing")

    # Send the response
    await update.message.reply_text(text, parse_mode="HTML")

# Start command handler
async def start(update: Update, context: CallbackContext):
    await update.message.reply_text("Hello! I am your chatbot. How can I assist you?")

# Message handler to reply with Gemini API response
async def handle_message(update: Update, context: CallbackContext):
    user_message = update.message.text

    # Show typing indicator
    await context.bot.send_chat_action(chat_id=update.message.chat_id, action="typing")

    # Check for custom responses first
    custom_response = get_custom_response(user_message)
    if custom_response:
        await send_response(update, context, custom_response)
        return

    # Check for custom commands (weather, news, etc.)
    if user_message.lower().startswith("weather in"):
        city = user_message[len("weather in"):].strip()
        response = await get_weather(city)
    elif user_message.lower() == "latest news":
        response = await get_news()
    elif user_message.lower().startswith("translate"):
        parts = user_message.split(" to ")
        if len(parts) == 2:
            text = parts[0][len("translate"):].strip()
            target_language = parts[1].strip()
            response = await translate_text(text, target_language)
        else:
            response = "Please provide a valid translation request."
    else:
        # Get the response from Gemini API
        try:
            if is_inappropriate(user_message):
                response = "Sorry, I cannot respond to that."
            else:
                response = model.generate_content(user_message).text
        except Exception as e:
            response = f"Error: {str(e)}"

    # Format the response (bold text)
    formatted_response = format_response(response)

    # Send the response
    await send_response(update, context, formatted_response)

# File handler for images, documents, etc.
async def handle_file(update: Update, context: CallbackContext):
    # Get the file object
    file = await update.message.document.get_file()
    file_path = await file.download_to_drive()

    # Get the file name and extension
    file_name = update.message.document.file_name
    file_extension = file_name.split(".")[-1].lower()

    # Show typing indicator
    await context.bot.send_chat_action(chat_id=update.message.chat_id, action="typing")

    # Process the file based on its type
    if file_extension in ["pdf", "docx", "pptx"]:
        if file_extension == "pdf":
            text = extract_text_from_pdf(file_path)
        elif file_extension == "docx":
            text = extract_text_from_docx(file_path)
        elif file_extension == "pptx":
            text = extract_text_from_pptx(file_path)
        
        if text:
            summary = summarize_text_with_gemini(text)
            response = summary if summary else "Failed to summarize the file."
        else:
            response = "Failed to extract text from the file."
    elif file_extension in ["png", "jpg", "jpeg"]:
        labels_with_scores = analyze_image_with_clarifai(file_path)
        if labels_with_scores:
            paragraph = generate_paragraph_with_gemini(labels_with_scores)
            response = paragraph if paragraph else "Failed to analyze the image."
        else:
            response = "Failed to analyze the image."
    else:
        response = "Unsupported file type."

    # Send the response
    await send_response(update, context, response)

# Main function to set up bot
def main():
    app = Application.builder().token(TELEGRAM_BOT_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_file))

    print("Bot is running...")
    app.run_polling()

if __name__ == "__main__":
    main()